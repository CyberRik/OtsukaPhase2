"""Text extraction per file type — the deterministic, GPU-free reading layer.

Each extractor returns plain text (truncated to `config.WORKSPACE_MAX_CHARS`) and
NEVER raises: a corrupt or unsupported file yields an empty string plus a note, so
one bad document degrades its own `extract` task without failing the run. Parsers
are imported lazily so a missing optional lib only affects that one format.
"""
from __future__ import annotations

from pathlib import Path

from senpai import config


def _cap(text: str) -> tuple[str, bool]:
    """Truncate to the char cap; return (text, truncated?)."""
    limit = config.WORKSPACE_MAX_CHARS
    if len(text) <= limit:
        return text, False
    return text[:limit].rstrip() + "\n… [truncated]", True


def _txt(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def _pdf(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    lines: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"[slide {i}]")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        lines.append(t)
    return "\n".join(lines)


def _xlsx(path: Path) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(str(path), read_only=True, data_only=True)
    lines: list[str] = []
    for ws in wb.worksheets:
        lines.append(f"[sheet {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append(" | ".join(cells))
            if len("\n".join(lines)) > config.WORKSPACE_MAX_CHARS:
                break  # stop early — a huge sheet can't help past the cap
    wb.close()
    return "\n".join(lines)


_EXTRACTORS = {
    ".txt": _txt, ".md": _txt,
    ".pdf": _pdf, ".docx": _docx, ".pptx": _pptx, ".xlsx": _xlsx,
}


def extract_text(path: Path) -> dict:
    """Read one document to text. Returns a structured dict (never raises):
    {text, chars, truncated, ext, error?}. Unknown/failed formats → empty text."""
    ext = path.suffix.lower()
    fn = _EXTRACTORS.get(ext)
    if fn is None:
        return {"text": "", "chars": 0, "truncated": False, "ext": ext,
                "error": f"unsupported extension {ext}"}
    try:
        raw = (fn(path) or "").strip()
    except Exception as e:  # noqa: BLE001 — a bad file degrades, never crashes
        return {"text": "", "chars": 0, "truncated": False, "ext": ext,
                "error": f"{type(e).__name__}: {e}"}
    text, truncated = _cap(raw)
    return {"text": text, "chars": len(text), "truncated": truncated, "ext": ext}
