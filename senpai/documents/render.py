"""Shared, LLM-free renderer — the only place python-pptx / python-docx are used.

Both the grounded tools (proposal/ringisho) and the general tools (pptx/docx) build
a normalized *spec* and hand it here. Keeping all rendering in one module means the
binary-format code is written and tested once.

Spec shapes
-----------
deck_spec (PPTX):
    {"slides": [
        {"layout": "title", "title": str, "subtitle": str},
        {"layout": "content", "title": str, "bullets": [str], "notes": str},
        ...
    ]}
doc_spec (DOCX):
    {"title": str, "subtitle": str,
     "sections": [{"heading": str, "body": [str]}, ...]}

`python-pptx` / `python-docx` are imported lazily so a missing lib can never break
the import of senpai.tools (mirrors senpai/tools/gcal.py).
"""
from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path

from senpai import config


def output_path(kind: str, slug: str, ext: str) -> Path:
    """A unique, safe path under GENERATED_DIR, e.g. proposal_D001_20260616-1430.pptx.
    Creates the dir on first use."""
    config.GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    safe = re.sub(r"[^A-Za-z0-9_-]+", "_", (slug or "doc")).strip("_") or "doc"
    stamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    return config.GENERATED_DIR / f"{kind}_{safe}_{stamp}.{ext}"


# Layouts we populate, mapped to the Otsuka template's layout names (and the
# blank default's, as a fallback). The template is hand-built: its "Title Only"
# layout (タイトルのみ) carries NO body placeholder, so the content layout MUST be
# resolved by name — blindly using slide_layouts[1] would silently drop bullets.
_TITLE_LAYOUT_NAMES = ("タイトル スライド", "Title Slide")
_CONTENT_LAYOUT_NAMES = ("タイトルとコンテンツ", "Title and Content")


def _layout(prs, names, fallback_idx):
    """Pick a slide layout by name (Otsuka template / Office default), else by index."""
    by_name = {layout.name: layout for layout in prs.slide_layouts}
    for name in names:
        if name in by_name:
            return by_name[name]
    return prs.slide_layouts[fallback_idx]


# Vector section icons — a single-glyph badge (colored circle + symbol), never an
# embedded image. Keeps every visual generated from code: nothing that could carry
# a stock photo's own context (a specific office, a specific person) into a deck
# about an unrelated customer. Otsuka blue / teal match the ribbon + bullet cards.
_OTSUKA_BLUE = (0x00, 0x55, 0xA4)
_TEAL = (0x14, 0xB8, 0xA6)
_AMBER = (0xE0, 0x7A, 0x2E)
_GRAY = (0x6B, 0x72, 0x80)
# The source deck's own convention (senpai/data/templates/otsuka_source.pptx):
# every number/price/percentage that appears in a body sentence is pulled out in
# bold navy, not left in the flat body color — it's how the real deck makes a
# grounded figure the thing your eye lands on. Reused here as text styling, not a
# new drawn shape.
_STAT_NAVY = (0x00, 0x20, 0x60)
_STAT_RE = re.compile(
    r"(¥[\d,]+(?:\.\d+)?|\d+(?:\.\d+)?\s?%|\d+(?:,\d{3})*(?:\.\d+)?"
    r"(?:万円|円|件|日間?|ヶ月|ケ月|か月|倍|名|社|回|年|時間))")


def _add_styled_runs(paragraph, text, size, base_color):
    """Split `text` on grounded numeric/price/percentage tokens and bold+navy
    them, leaving the rest at `base_color` — matches the source template's own
    pull-the-number-out convention instead of one flat run."""
    from pptx.dml.color import RGBColor

    parts = _STAT_RE.split(text)
    for i, part in enumerate(parts):
        if not part:
            continue
        run = paragraph.add_run()
        run.text = part
        run.font.size = size
        if i % 2 == 1:  # a captured stat token
            run.font.bold = True
            run.font.color.rgb = RGBColor(*_STAT_NAVY)
        else:
            run.font.color.rgb = base_color
_ICONS = {
    "challenge": ("!", _AMBER),
    "solution": ("✓", _TEAL),
    "roi": ("↑", _OTSUKA_BLUE),
    "next": ("→", _OTSUKA_BLUE),
    "schedule": ("▤", _GRAY),
    "background": ("★", _TEAL),
    "assessment": ("⚙", _GRAY),
    "summary": ("§", _OTSUKA_BLUE),
}


def _render_comparison_png(categories: list[str], values: list[float],
                           value_labels: list[str] | None = None) -> "io.BytesIO":
    """A styled horizontal before/after bar comparison, rendered server-side with
    matplotlib and handed back as PNG bytes for `add_picture` — real numbers in,
    a custom-designed image out (no native-chart legend/gridline styling that
    doesn't match the deck's flat, no-shadow look). Deterministic: same numbers
    always draw the same image, nothing invented or sampled."""
    import io

    import matplotlib
    matplotlib.use("Agg")  # headless — no display/GUI backend needed on a server
    import matplotlib.pyplot as plt
    import matplotlib.font_manager as fm

    for candidate in ("Yu Gothic", "Meiryo", "MS Gothic"):
        if any(candidate.lower() in f.name.lower() for f in fm.fontManager.ttflist):
            plt.rcParams["font.family"] = candidate
            break

    navy = "#00205A"       # matches _STAT_NAVY
    gray = "#94A3B8"
    accent = "#E13365"     # the source template's own punchline-pink

    fig, ax = plt.subplots(figsize=(7.5, 2.6), dpi=200)
    bar_colors = [gray, navy] if len(values) > 1 else [navy]
    bars = ax.barh(categories, values, color=bar_colors, height=0.5)

    labels = value_labels or [f"¥{v:,.0f}" for v in values]
    for bar, label in zip(bars, labels):
        ax.text(bar.get_width() + max(values) * 0.02, bar.get_y() + bar.get_height() / 2,
               label, va="center", ha="left", fontsize=13, fontweight="bold", color=navy)

    if len(values) == 2 and values[0]:
        pct = round((1 - values[1] / values[0]) * 100)
        if pct > 0:
            ax.text(max(values) * 0.5, -0.75, f"▼ {pct}%",
                   fontsize=15, fontweight="bold", color=accent, ha="center")

    ax.set_xlim(0, max(values) * 1.35)
    ax.invert_yaxis()
    for spine in ax.spines.values():
        spine.set_visible(False)
    ax.set_xticks([])
    ax.tick_params(axis="y", length=0, labelsize=13)
    fig.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format="png", transparent=True)
    plt.close(fig)
    buf.seek(0)
    return buf


def _add_icon_badge(slide, icon_key, x, y, size):
    """A small colored-circle glyph badge — decorative theme marker, not data."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    glyph, color = _ICONS.get(icon_key, ("•", _OTSUKA_BLUE))
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    badge.fill.solid()
    badge.fill.fore_color.rgb = RGBColor(*color)
    badge.line.fill.background()
    badge.shadow.inherit = False
    tf = badge.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = glyph
    run.font.size = Pt(20)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    return badge


def render_pptx(deck_spec: dict, path: Path) -> Path:
    """Render a deck spec to a .pptx file at `path`. Returns the path.

    Opens the committed Otsuka brand template (config.PPTX_TEMPLATE_PATH) as the
    base so the deck inherits its masters/layouts/theme; falls back to python-pptx's
    blank default if the template is missing (e.g. in CI without the asset).
    """
    from pptx import Presentation
    from pptx.util import Pt, Inches
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.chart.data import CategoryChartData
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    tmpl = config.PPTX_TEMPLATE_PATH
    prs = Presentation(str(tmpl)) if tmpl.exists() else Presentation()
    title_layout = _layout(prs, _TITLE_LAYOUT_NAMES, 0)       # Title Slide
    content_layout = _layout(prs, _CONTENT_LAYOUT_NAMES, 1)   # Title and Content

    slides = deck_spec.get("slides") or []
    if not slides:  # never produce an empty deck
        slides = [{"layout": "title", "title": deck_spec.get("title", "Document"),
                   "subtitle": deck_spec.get("subtitle", "")}]

    for spec in slides:
        is_title = spec.get("layout") == "title"
        slide = prs.slides.add_slide(title_layout if is_title else content_layout)



        if slide.shapes.title is not None:
            slide.shapes.title.text = str(spec.get("title", ""))

        icon_key = spec.get("icon")
        if icon_key and not is_title:
            _add_icon_badge(slide, icon_key, prs.slide_width - Inches(1.0),
                            Inches(0.3), Inches(0.55))

        if is_title:
            # subtitle placeholder (idx 1) on the title layout
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    ph.text = str(spec.get("subtitle", ""))
                    break
            continue

        is_table = spec.get("layout") == "table"
        is_chart = spec.get("layout") == "chart"

        if is_table or is_chart:
            # Remove the body placeholder to avoid "Click to add text" prompt
            for shape in list(slide.shapes):
                if shape.is_placeholder and shape.placeholder_format.idx == 1:
                    sp = shape._sp
                    sp.getparent().remove(sp)
            
            # Use typical positioning for content
            x, y, cx, cy = Inches(1), Inches(1.5), Inches(8), Inches(4.5)

            if is_table:
                table_data = spec.get("table", {})
                headers = table_data.get("headers", [])
                rows = table_data.get("rows", [])
                num_rows = len(rows) + (1 if headers else 0)
                num_cols = max(len(headers), max((len(r) for r in rows), default=0)) if headers or rows else 1
                
                if num_rows > 0 and num_cols > 0:
                    table_shape = slide.shapes.add_table(num_rows, num_cols, x, y, cx, cy)
                    table = table_shape.table
                    
                    row_idx = 0
                    if headers:
                        for col_idx, header in enumerate(headers):
                            cell = table.cell(row_idx, col_idx)
                            cell.text = str(header)
                            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                            # Make header bold
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.bold = True
                        row_idx += 1
                    
                    for row in rows:
                        for col_idx, item in enumerate(row):
                            if col_idx < num_cols:
                                cell = table.cell(row_idx, col_idx)
                                cell.text = str(item)
                                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                        row_idx += 1

            elif is_chart and spec.get("chart", {}).get("renderer") == "mpl":
                # A custom-designed comparison, rendered server-side (matplotlib)
                # and embedded as a picture — real numbers in, brand-matched
                # styling out, no native-chart legend/gridlines to fight with.
                cs = spec["chart"]
                categories = cs.get("categories", [])
                values = (cs.get("series") or [{}])[0].get("values", [])
                labels = cs.get("value_labels") or [f"¥{v:,.0f}" for v in values]
                png = _render_comparison_png(categories, values, labels)
                slide.shapes.add_picture(png, Inches(1), Inches(2.0), width=Inches(8))

                # The image is illustrative — the actual figures also need to
                # exist as real text (accessibility, search, copy-paste), not
                # only as pixels. A compact caption line restates them, with the
                # same number-highlighting every other slide uses.
                caption = slide.shapes.add_textbox(Inches(1), Inches(5.1), Inches(8), Inches(0.5))
                cap_tf = caption.text_frame
                cap_tf.word_wrap = True
                cap_p = cap_tf.paragraphs[0]
                cap_p.alignment = PP_ALIGN.CENTER
                _add_styled_runs(cap_p, "　|　".join(f"{c}: {v}" for c, v in zip(categories, labels)),
                                 Pt(12), RGBColor(*_GRAY))

            elif is_chart:
                chart_data_spec = spec.get("chart", {})
                chart_data = CategoryChartData()
                chart_data.categories = chart_data_spec.get("categories", [])
                for series in chart_data_spec.get("series", []):
                    chart_data.add_series(series.get("name", ""), series.get("values", []))

                # "type": "doughnut" for a single-metric proportion (e.g. discount
                # rate) instead of the default column-bar comparison.
                chart_type = (XL_CHART_TYPE.DOUGHNUT if chart_data_spec.get("type") == "doughnut"
                             else XL_CHART_TYPE.COLUMN_CLUSTERED)
                if chart_type == XL_CHART_TYPE.DOUGHNUT:
                    cx = cy = Inches(3.5)  # a doughnut reads best square, not stretched
                chart = slide.shapes.add_chart(
                    chart_type, x, y, cx, cy, chart_data
                ).chart

                # Add data labels
                plot = chart.plots[0]
                plot.has_data_labels = True
                for series in plot.series:
                    for point in series.points:
                        point.data_label.has_text_frame = True

            # If notes exist, add them
            notes = str(spec.get("notes", "") or "").strip()
            if notes:
                slide.notes_slide.notes_text_frame.text = notes
            continue

        is_timeline = spec.get("layout") == "timeline"
        if is_timeline:
            # A visual process flow (connected boxes + arrows) instead of a plain
            # table — same deterministic phase data, more legible for a rep to
            # walk a customer through in a meeting.
            for shape in list(slide.shapes):
                if shape.is_placeholder and shape.placeholder_format.idx == 1:
                    sp = shape._sp
                    sp.getparent().remove(sp)

            phases = spec.get("phases") or []
            n = len(phases)
            if n:
                margin = Inches(0.6)
                arrow_w = Inches(0.35)
                box_w = int((prs.slide_width - 2 * margin - (n - 1) * arrow_w) / n)
                box_h = Inches(1.6)
                y = Inches(2.8)
                x = margin
                for i, ph in enumerate(phases):
                    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
                    box.fill.solid()
                    box.fill.fore_color.rgb = (RGBColor(0xF0, 0xF4, 0xF8) if i % 2 == 0
                                               else RGBColor(0xE2, 0xEC, 0xF5))
                    box.line.color.rgb = RGBColor(*_OTSUKA_BLUE)
                    box.shadow.inherit = False
                    tf = box.text_frame
                    tf.word_wrap = True
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    p0 = tf.paragraphs[0]
                    p0.alignment = PP_ALIGN.CENTER
                    r0 = p0.add_run()
                    r0.text = str(ph.get("label", ""))
                    r0.font.bold = True
                    r0.font.size = Pt(13)
                    r0.font.color.rgb = RGBColor(0, 0, 0)
                    p1 = tf.add_paragraph()
                    p1.alignment = PP_ALIGN.CENTER
                    _add_styled_runs(p1, str(ph.get("duration", "")), Pt(11), RGBColor(*_GRAY))
                    detail = str(ph.get("detail", "")).strip()
                    if detail:
                        p2 = tf.add_paragraph()
                        p2.alignment = PP_ALIGN.CENTER
                        r2 = p2.add_run()
                        r2.text = detail
                        r2.font.size = Pt(9)
                        r2.font.color.rgb = RGBColor(*_GRAY)
                    x += box_w
                    if i < n - 1:
                        arrow = slide.shapes.add_shape(
                            MSO_SHAPE.RIGHT_ARROW, x, y + box_h // 2 - Inches(0.15),
                            arrow_w, Inches(0.3))
                        arrow.fill.solid()
                        arrow.fill.fore_color.rgb = RGBColor(*_OTSUKA_BLUE)
                        arrow.line.fill.background()
                        arrow.shadow.inherit = False
                        x += arrow_w

            notes = str(spec.get("notes", "") or "").strip()
            if notes:
                slide.notes_slide.notes_text_frame.text = notes
            continue

        bullets = [str(b) for b in (spec.get("bullets") or []) if str(b).strip()]

        if bullets and not (is_table or is_chart):
            # Remove the default body placeholder so we can draw custom shapes
            for shape in list(slide.shapes):
                if shape.is_placeholder and shape.placeholder_format.idx == 1:
                    sp = shape._sp
                    sp.getparent().remove(sp)

            # Left-aligned accent-stripe cards — no border, no drop shadow (the
            # centered, shadowed "pill button" look read as a default-autoshape
            # placeholder, not a designed slide). A thin color bar + a small vector
            # dot substitute for the border, and the block is vertically centered
            # in the available area instead of always starting at a fixed y and
            # leaving the rest of the slide empty for short bullet lists.
            margin = Inches(0.8)
            width = prs.slide_width - 2 * margin
            # Roughly two lines' worth of height once text wraps past ~44 chars
            # (word_wrap is on regardless — this only sizes the card, not the text).
            heights = [Inches(1.0) if len(b) > 44 else Inches(0.62) for b in bullets]
            gap = Inches(0.22)
            content_top, content_bottom = Inches(1.95), Inches(6.9)
            total_h = sum(heights, start=0) + gap * (len(bullets) - 1)
            y = content_top + max(0, (content_bottom - content_top - total_h)) // 2
            x = margin

            for b, h in zip(bullets, heights):
                card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, h)
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
                card.line.fill.background()
                card.shadow.inherit = False
                card.adjustments[0] = 0.06  # subtler corner radius than the default

                accent = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(0.07), h)
                accent.fill.solid()
                accent.fill.fore_color.rgb = RGBColor(*_OTSUKA_BLUE)
                accent.line.fill.background()
                accent.shadow.inherit = False

                dot_size = Inches(0.14)
                dot = slide.shapes.add_shape(
                    MSO_SHAPE.OVAL, x + Inches(0.3), y + h // 2 - dot_size // 2, dot_size, dot_size)
                dot.fill.solid()
                dot.fill.fore_color.rgb = RGBColor(*_TEAL)
                dot.line.fill.background()
                dot.shadow.inherit = False

                tf = card.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.margin_left = Inches(0.6)
                tf.margin_right = Inches(0.3)
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                _add_styled_runs(p, b, Pt(16), RGBColor(0x1F, 0x29, 0x37))

                y += h + gap

        notes = str(spec.get("notes", "") or "").strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(str(path))
    return path


def render_docx(doc_spec: dict, path: Path) -> Path:
    """Render a doc spec to a .docx file at `path`. Returns the path."""
    from docx import Document

    doc = Document()
    title = str(doc_spec.get("title", "") or "")
    if title:
        doc.add_heading(title, level=0)
    subtitle = str(doc_spec.get("subtitle", "") or "").strip()
    if subtitle:
        doc.add_paragraph(subtitle).italic = True

    for section in doc_spec.get("sections") or []:
        heading = str(section.get("heading", "") or "").strip()
        if heading:
            doc.add_heading(heading, level=1)
        for para in section.get("body") or []:
            text = str(para).strip()
            if not text:
                continue
            # A leading "- " or "・" renders as a bullet list item.
            if text[:2] in ("- ", "・") or text.startswith("• "):
                doc.add_paragraph(text.lstrip("-・• ").strip(), style="List Bullet")
            else:
                doc.add_paragraph(text)

    doc.save(str(path))
    return path
