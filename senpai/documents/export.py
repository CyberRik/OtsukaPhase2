"""Export an HTML deck (senpai/documents/html_render.py) to PDF and editable PPTX.

Both exports drive the same headless Chromium that senpai/tools/crawl.py already uses
(Playwright), so no new runtime stack. Everything here is synchronous and expects to be
called from the sync tool layer (the same context crawl_site runs in); if Chromium is
unavailable it degrades gracefully — `export_html_deck` returns which artifacts it managed
to produce and the caller falls back to the native python-pptx renderer.

PDF path: `page.pdf()` of the print-media HTML — pixel-perfect, one slide per page.

PPTX path (the "editable text over a faithful background" strategy):
  1. Measure every `.pptx-text` element's box + computed font (via page.evaluate).
  2. Hide those text elements and screenshot each `.slide` → a full-bleed background that
     carries all the decoration (accent bars, cards, charts, table rules) but no text.
  3. Build a 16:9 PPTX: each slide = the background picture + one NATIVE python-pptx
     text box per measured element, positioned by px→EMU conversion, with the figure
     tokens re-bolded in navy (render._STAT_RE convention). Result looks ~identical to the
     HTML yet every word is real, selectable and editable in PowerPoint.
"""
from __future__ import annotations

import io
import re
from pathlib import Path

from senpai.documents.html_render import SLIDE_W, SLIDE_H, _STAT_RE

# Standard PowerPoint 16:9 canvas = 13.333in × 7.5in. Our HTML slide is SLIDE_W×SLIDE_H
# CSS px at 96 dpi, i.e. exactly that canvas, so 1 px = 9525 EMU and 1 px = 0.75 pt.
_EMU_PER_PX = 9525
_PT_PER_PX = 0.75
_SLIDE_W_EMU = SLIDE_W * _EMU_PER_PX   # 12,192,000
_SLIDE_H_EMU = SLIDE_H * _EMU_PER_PX   # 6,858,000

# A JP-friendly typeface that ships with Windows/Mac PowerPoint, so the editable text
# renders correctly on the user's machine even though Chromium laid it out in Noto.
_PPTX_FONT = "Yu Gothic"

# One evaluate() over the whole document: per slide, its notes + every editable text run's
# box (relative to the slide) and computed style.
_MEASURE_JS = r"""
() => Array.from(document.querySelectorAll('.slide')).map(slide => {
  const sr = slide.getBoundingClientRect();
  const texts = Array.from(slide.querySelectorAll('.pptx-text')).map(el => {
    const r = el.getBoundingClientRect();
    const cs = getComputedStyle(el);
    return {
      x: r.x - sr.x, y: r.y - sr.y, w: r.width, h: r.height,
      text: (el.innerText || '').replace(/ /g, ' ').trim(),
      size: parseFloat(cs.fontSize) || 16,
      weight: parseInt(cs.fontWeight, 10) || 400,
      italic: cs.fontStyle === 'italic',
      color: cs.color || 'rgb(26,35,48)',
      align: cs.textAlign || 'left',
      nowrap: cs.whiteSpace === 'nowrap' || cs.whiteSpace === 'pre'
    };
  }).filter(t => t.text.length > 0);
  return { notes: slide.getAttribute('data-notes') || '', texts };
});
"""

# Injected before screenshotting: make the editable text invisible in the background so
# it isn't baked in twice, WITHOUT disturbing layout or decoration. `color: transparent`
# (not visibility:hidden) is deliberate — visibility:hidden would also drop the element's
# own background/border, which would erase table-header fills and card backgrounds that
# live on the same element as the text. The `*` clause also neutralizes the navy figure
# runs so they don't survive as baked color under the overlaid text.
_EXPORT_CSS = """
.pptx-text, .pptx-text * { color: transparent !important; text-shadow: none !important; }
.slide { box-shadow: none !important; border-radius: 0 !important; }
.nav { display: none !important; }
"""


def render_deck(spec: dict, *, kind: str, slug: str, lang: str = "ja") -> dict[str, "Path"]:
    """Render a deck spec to HTML + PDF + editable PPTX and return the produced files
    as {"pptx": Path, "html": Path, "pdf": Path?}. This is the one HTML-first path shared
    by every deck generator (generate_pptx, generate_proposal, the planner). The PPTX is
    always produced: if headless Chromium is unavailable the browser exports are skipped
    and it falls back to the native python-pptx renderer, so a deck is never lost.

    `kind` names the file family (e.g. 'pptx', 'proposal'); `slug` is the human filename
    stem (deal id or title)."""
    from senpai.documents import html_render
    from senpai.documents.render import output_path, render_pptx

    html = html_render.render_html(spec, lang=lang)
    html_path = output_path(kind, slug, "html")
    html_path.write_text(html, encoding="utf-8")
    pptx_path = output_path(kind, slug, "pptx")
    pdf_path = output_path(kind, slug, "pdf")

    produced = export_html_deck(html, pptx_path=pptx_path, pdf_path=pdf_path)
    if not produced.get("pptx"):
        render_pptx(spec, pptx_path)  # browser unavailable → native fallback

    files: dict[str, Path] = {"pptx": pptx_path, "html": html_path}
    if produced.get("pdf"):
        files["pdf"] = pdf_path
    return files


def export_html_deck(html: str, *, pptx_path: Path | None = None,
                     pdf_path: Path | None = None) -> dict[str, bool]:
    """Produce the requested artifacts from `html`. Returns e.g. {"pptx": True,
    "pdf": True}. On any browser failure returns {} so the caller can fall back."""
    try:
        from playwright.sync_api import sync_playwright  # type: ignore
    except Exception:
        return {}

    done: dict[str, bool] = {}
    pw = browser = None
    try:
        pw = sync_playwright().start()
        browser = pw.chromium.launch(headless=True)
        ctx = browser.new_context(
            viewport={"width": SLIDE_W, "height": SLIDE_H}, device_scale_factor=2)
        page = ctx.new_page()
        page.set_content(html, wait_until="load")
        try:
            page.evaluate("() => document.fonts && document.fonts.ready")
        except Exception:
            pass
        page.wait_for_timeout(180)  # let fonts/layout settle

        if pdf_path is not None:
            try:
                page.emulate_media(media="print")
                page.pdf(path=str(pdf_path), width=f"{SLIDE_W}px", height=f"{SLIDE_H}px",
                         print_background=True, margin={"top": "0", "right": "0",
                                                        "bottom": "0", "left": "0"})
                page.emulate_media(media="screen")
                done["pdf"] = True
            except Exception:
                done["pdf"] = False

        if pptx_path is not None:
            try:
                _render_pptx_from_page(page, pptx_path)
                done["pptx"] = True
            except Exception:
                done["pptx"] = False
    except Exception:
        return done
    finally:
        for obj, meth in ((browser, "close"), (pw, "stop")):
            try:
                if obj is not None:
                    getattr(obj, meth)()
            except Exception:
                pass
    return done


def _render_pptx_from_page(page, pptx_path: Path) -> None:
    """Measure text, screenshot text-hidden slide backgrounds, assemble the PPTX."""
    from pptx import Presentation
    from pptx.util import Emu, Pt
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

    metas = page.evaluate(_MEASURE_JS)
    handles = page.query_selector_all(".slide")
    page.add_style_tag(content=_EXPORT_CSS)  # hide text + flatten AFTER measuring
    backgrounds = [h.screenshot(type="png") for h in handles]

    align_map = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT,
                 "justify": PP_ALIGN.JUSTIFY, "left": PP_ALIGN.LEFT, "start": PP_ALIGN.LEFT}

    prs = Presentation()
    prs.slide_width = Emu(_SLIDE_W_EMU)
    prs.slide_height = Emu(_SLIDE_H_EMU)
    blank = prs.slide_layouts[6]

    for meta, bg in zip(metas, backgrounds):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(io.BytesIO(bg), 0, 0,
                                 width=Emu(_SLIDE_W_EMU), height=Emu(_SLIDE_H_EMU))
        for t in meta.get("texts", []):
            # Give each box a little horizontal slack so PowerPoint/LibreOffice metrics
            # (Yu Gothic) — slightly wider than Chromium's Noto — don't wrap a line that
            # fit in the HTML. Extend symmetrically so centered text stays centered, and
            # clamp x so we never push off the left edge.
            slack = max(14, round(t["w"] * 0.06))
            x = max(0, int((t["x"] - slack / 2) * _EMU_PER_PX))
            box = slide.shapes.add_textbox(
                Emu(x), Emu(int(t["y"] * _EMU_PER_PX)),
                Emu(int((t["w"] + slack) * _EMU_PER_PX)), Emu(int(max(t["h"], 8) * _EMU_PER_PX)))
            tf = box.text_frame
            # Honor the HTML's white-space: a nowrap element (e.g. a big stat figure) must
            # not wrap in PowerPoint — it overflows its box symmetrically instead, which for
            # centered text keeps it centered rather than breaking onto a second line.
            tf.word_wrap = not t.get("nowrap")
            for side in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
                setattr(tf, side, Emu(0))
            try:
                tf.auto_size = MSO_AUTO_SIZE.NONE
            except Exception:
                pass
            align = align_map.get((t.get("align") or "left").lower(), PP_ALIGN.LEFT)
            size = Pt(t["size"] * _PT_PER_PX)
            color = _parse_rgb(t["color"])
            bold = t.get("weight", 400) >= 600
            italic = bool(t.get("italic"))
            # A single element may hold several visual lines (e.g. a cover subtitle with
            # line breaks) — one PPTX paragraph per line so they stay editable and stacked.
            for li, line in enumerate(t["text"].split("\n")):
                para = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
                para.alignment = align
                _add_runs(para, line, size, color, bold=bold, italic=italic)
        notes = (meta.get("notes") or "").strip()
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    prs.save(str(pptx_path))


def _add_runs(paragraph, text: str, size, base_color, *, bold: bool, italic: bool) -> None:
    """Emit runs for `text`, bolding grounded figure tokens in navy (same rule as
    render._add_styled_runs) while keeping the element's base weight/italic/color."""
    from pptx.dml.color import RGBColor
    from pptx.oxml.ns import qn

    _STAT_NAVY = RGBColor(0x00, 0x20, 0x60)
    parts = _STAT_RE.split(text or "")
    for i, part in enumerate(parts):
        if not part:
            continue
        run = paragraph.add_run()
        run.text = part
        run.font.size = size
        run.font.italic = italic
        is_fig = i % 2 == 1
        run.font.bold = True if is_fig else bold
        run.font.color.rgb = _STAT_NAVY if is_fig else base_color
        # Set the typeface for Latin, East-Asian and complex scripts so Japanese renders
        # in the intended font rather than the theme default.
        run.font.name = _PPTX_FONT
        rpr = run._r.get_or_add_rPr()
        for tag in ("a:latin", "a:ea", "a:cs"):
            el = rpr.find(qn(tag))
            if el is None:
                el = rpr.makeelement(qn(tag), {})
                rpr.append(el)
            el.set("typeface", _PPTX_FONT)


_RGB_RE = re.compile(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)")


def _parse_rgb(css: str):
    from pptx.dml.color import RGBColor
    m = _RGB_RE.match(css or "")
    if m:
        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    return RGBColor(0x1A, 0x23, 0x30)


_HEADING_RE = re.compile(r"^#{1,6}\s+(.*)$")
_RULE_RE = re.compile(r"^-{3,}$")
_BULLET_RE = re.compile(r"^[-*]\s+(.*)$")


def _text_to_doc_spec(text: str, title: str) -> dict:
    sections: list[dict] = []
    current = {"heading": "", "body": []}

    def flush() -> None:
        if current["body"] or current["heading"]:
            sections.append({"heading": current["heading"], "body": list(current["body"])})

    for raw_line in (text or "").replace("\r", "").split("\n"):
        line = raw_line.strip()
        if not line or _RULE_RE.match(line):
            continue
        heading = _HEADING_RE.match(line)
        if heading:
            flush()
            current = {"heading": heading.group(1).strip(), "body": []}
            continue
        bullet = _BULLET_RE.match(line)
        current["body"].append(f"- {bullet.group(1).strip()}" if bullet else line)
    flush()

    if not sections:
        sections = [{"heading": "", "body": [(text or "").strip() or "(no content)"]}]

    return {"title": title or "Export", "subtitle": "", "sections": sections}


def export_text_as_docx(text: str, title: str = "", slug: str = "") -> dict:
    """Render `text` verbatim (parsed, not LLM-rewritten) to a .docx and register
    it for download. Returns the registry record (doc_id, filename, download_url)."""
    from senpai.documents import registry, render
    spec = _text_to_doc_spec(text, title)
    path = render.output_path("export", slug or title or "chat", "docx")
    render.render_docx(spec, path)
    return registry.register("export", path)

