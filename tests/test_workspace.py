"""Workspace capability — sandboxed local document access on the orchestration engine.

Covers the four milestone-1 guarantees: (1) the sandbox rejects path escapes,
(2) every declared file type extracts, (3) `find` fans out at runtime into one
`extract` task per document (the first production use of `ctx.expand`), and
(4) the gather reduces to grounded evidence with file citations — degrading, never
crashing, on a missing/empty workspace. No GPU/model needed.
"""
from __future__ import annotations

import pytest
from docx import Document as Docx
from openpyxl import Workbook
from pptx import Presentation

from senpai import config
from senpai.workspace import sandbox
from senpai.workspace.capabilities import build_registry
from senpai.workspace.gather import workspace_evidence
from senpai.workspace.plan import workspace_plan


@pytest.fixture
def ws(tmp_path, monkeypatch):
    """A populated sandbox root pointed at a tmp dir (one file per supported type)."""
    monkeypatch.setattr(config, "WORKSPACE_ROOT", tmp_path.resolve())
    (tmp_path / "notes.txt").write_text("Endo Kogyo サーバー更改の提案メモ。", encoding="utf-8")
    (tmp_path / "summary.md").write_text("# 要約\n決裁者は情報システム部長。", encoding="utf-8")
    d = Docx(); d.add_paragraph("Endo Kogyo との商談記録。"); d.save(str(tmp_path / "memo.docx"))
    p = Presentation()
    s = p.slides.add_slide(p.slide_layouts[5])
    s.shapes.title.text = "提案 Endo Kogyo"
    p.save(str(tmp_path / "deck.pptx"))
    wb = Workbook(); wsheet = wb.active; wsheet.append(["item", "yen"]); wsheet.append(["NAS08", 240000])
    wb.save(str(tmp_path / "quote.xlsx"))
    # a decoy outside the allowed set + a nested dir
    (tmp_path / "ignore.log").write_text("nope", encoding="utf-8")
    (tmp_path / "sub").mkdir()
    (tmp_path / "sub" / "deep.txt").write_text("nested Endo file", encoding="utf-8")
    return tmp_path


# --- sandbox ------------------------------------------------------------------
def test_sandbox_blocks_escape(ws):
    for bad in ("../../etc/passwd", "/etc/passwd", "../outside.txt"):
        with pytest.raises(sandbox.SandboxError):
            sandbox.safe_path(bad)


def test_sandbox_lists_only_allowed(ws):
    names = {p.name for p in sandbox.list_documents()}
    assert "notes.txt" in names and "deep.txt" in names   # recursive
    assert "ignore.log" not in names                       # extension filter


def test_missing_workspace_degrades(tmp_path, monkeypatch):
    monkeypatch.setattr(config, "WORKSPACE_ROOT", (tmp_path / "does-not-exist").resolve())
    assert sandbox.list_documents() == []
    res = workspace_evidence("anything")
    assert res["available"] == 0 and res["documents"] == []


# --- extraction (every declared type) ----------------------------------------
def test_extracts_every_type(ws):
    res = workspace_evidence("", limit=20)  # empty query → all, recency order
    by_ext = {d["ext"] for d in res["documents"]}
    assert {".txt", ".md", ".docx", ".pptx", ".xlsx"} <= by_ext
    text = {d["name"]: d["text"] for d in res["documents"]}
    assert "サーバー更改" in text["notes.txt"]
    assert "商談記録" in text["memo.docx"]
    assert "Endo Kogyo" in text["deck.pptx"]
    assert "240000" in text["quote.xlsx"]


# --- runtime fan-out (ctx.expand) --------------------------------------------
def test_find_fans_out_into_extract_tasks(ws):
    from senpai.orchestration import ExecutionEngine
    bundle = ExecutionEngine(build_registry()).run(workspace_plan("Endo"), lambda *_a, **_k: None)
    # one find task + one extract task per matched document
    find = bundle.get("find")
    assert find is not None and find.status == "ok"
    extract_ids = [tid for tid in bundle.fragments if tid.startswith("find:extract:")]
    assert len(extract_ids) == find.data["count"] >= 1     # DAG grew at runtime
    assert all(bundle.get(tid).capability == "workspace" for tid in extract_ids)


def test_find_relevance_and_citations(ws):
    res = workspace_evidence("Endo")
    assert res["documents"], "query should match the Endo files"
    assert all(c.startswith("file://") for c in res["citations"])
    # matched files carry 'Endo' in name OR were selected — never the .log decoy
    assert all(d["ext"] in config.WORKSPACE_EXTS for d in res["documents"])


def test_fanout_capped(ws, monkeypatch):
    monkeypatch.setattr(config, "WORKSPACE_MAX_FILES", 2)
    res = workspace_evidence("")     # empty query would otherwise take all
    assert len(res["documents"]) <= 2
